import base64
from pymupdf4llm import to_markdown
from docx import Document
from pptx import Presentation
import pandas as pd
from bs4 import BeautifulSoup
import json

class FileHandler:
    def __init__(self, file):
        """
        Initialize the FileHandler with a file object.

        :param file: The file object with attributes `mime` and `path`.
        """
        self.file = file
        self.content = None
        self.category = None

    def process_file(self):
        """
        Process the file based on its MIME type.

        :return: Processed content (encoded or converted based on the file type).
        """
        try:
            # Determine the category of the file
            self.category = self.categorize_mime(self.file.mime)
            print(self.category)

            handlers = {
                "image/jpeg": self._process_image,
                "image/png": self._process_image,
                "text/html": self._process_html,
                "text/plain": self._process_txt,
                "text/csv": self._process_spreadsheet,
                "application/octet-stream": self._process_binary,
                "application/json": self._process_json,
                "application/pdf": self._process_pdf,
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document": self._process_word,  # .docx
                "application/vnd.openxmlformats-officedocument.presentationml.presentation": self._process_ppt,  # .pptx
                "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": self._process_spreadsheet,  # .xlsx
            }

            handler = handlers.get(self.file.mime)
            if handler:
                handler()  # Call the corresponding function for the file type
            else:
                raise ValueError(f"Unsupported MIME type: {self.file.mime}")

            # Save the processed content along with the category
            self.content = {
                "file_name": self.file.name,
                "category": self.category,
                "data": self.content, 
            }
            return self.content

        except Exception as e:
            print(f"Error processing file: {e}")
            return None

    @staticmethod
    def categorize_mime(mime_type):
        """
        Categorize a file based on its MIME type.

        :param mime_type: The MIME type of the file.
        :return: Category of the file.
        """
        categories = {
            "image": ["image/jpeg", "image/png"],
            "code": ["application/octet-stream", "application/json"],
            "pandas": ["text/csv", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"],
            
            "document": [
                "text/plain", "application/pdf",
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                "application/vnd.openxmlformats-officedocument.presentationml.presentation"],
        }

        for category, mime_list in categories.items():
            if mime_type in mime_list:
                return category
        return "unknown"  

    def _process_image(self):
        """Encode image content to Base64."""
        with open(self.file.path, "rb") as f:
            b64 = base64.b64encode(f.read()).decode("utf-8")
        self.content = f"data:{self.file.mime};base64,{b64}" 


    def _process_html(self):
        """Extract text content from an HTML file."""
        with open(self.file.path, "r", encoding="utf-8") as f:
            self.content = BeautifulSoup(f.read(), "html.parser").get_text()

    def _process_txt(self):
        """Read and decode text file content."""
        with open(self.file.path, "r", encoding="utf-8") as f:
            self.content = f.read()

    def _process_binary(self):
        """Read and decode binary file content."""
        with open(self.file.path, "rb") as f:
            self.content = f.read().decode("utf-8")

    def _process_json(self):
        """
        Load and process content from a JSON file.
        """
        with open(self.file.path, "r", encoding="utf-8") as f:
                self.content = json.load(f)  

    def _process_pdf(self):
        """Get PDF content."""
        self.content = to_markdown(self.file.path)
        
    def _process_word(self):
        """Get content from a Word (.docx) file."""
        self.content = [
            para.text.strip()
            for para in Document(self.file.path).paragraphs
            if para.text.strip()
        ]

    def _process_ppt(self):
        """Get content from a PowerPoint (.pptx) file."""
        presentation = Presentation(self.file.path)
        self.content = [
            "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text"))
            for slide in presentation.slides
        ]

    def _process_spreadsheet(self):
        """Read content from a CSV or Excel file."""
        read_func = pd.read_csv if self.file.mime == "text/csv" else pd.read_excel
        df = read_func(self.file.path)
        self.content = df.to_csv(index=False)


# Example usage
# my_files = [file for file in input_message.elements]
# handler = FileHandler(my_files[0])
# processed_content = handler.process_file()
# print(processed_content)
